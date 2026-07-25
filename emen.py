from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

pr = Presentation()

# Set slide width and height for better fit (widescreen)
pr.slide_width = Inches(13.333)
pr.slide_height = Inches(7.5)

def add_content_slide(title, content_lines, font_size=18):
    slide_layout = pr.slide_layouts[1]
    slide = pr.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    for line in content_lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.space_after = Pt(4)
        if line.startswith("•") or line.startswith("✓") or line.startswith("✗") or line.startswith("1.") or line.startswith("2.") or line.startswith("3.") or line.startswith("4.") or line.startswith("5.") or line.startswith("6.") or line.startswith("7."):
            p.level = 0
        elif line.startswith("   -") or line.startswith("   o") or line.startswith("   •"):
            p.level = 1
            p.font.size = Pt(font_size - 2)
    
    return slide

# SLIDE 1: COVER PAGE
slide_layout = pr.slide_layouts[0]
slide = pr.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Computer Vision\nHow It Works, Advantages & Disadvantages"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
subtitle = slide.placeholders[1]
subtitle.text = """Based on: Introduction to Emerging Technology (Addis Ababa University, 2019)

Presented by:
No.  Student ID      Full Name
1    RU1066/18      EYEL FITSUM YIMAM
2    RU1069/18      EYEL TILAHUN ENDALE
3    RU3223/18      FADUMO MOHAMUD ABDALE
4    RU1105/18      FEBEN ASHENAFI SERBA
5    RU1106/18      FEBEN NEGASH DABELU
6    RU1115/18      FENET BIRIHANU MINTEFA
7    RU1187/18      FIRAOL TADESA HIRHA
8    RU1254/18      GAMACHU NAGASA MAGARSA
9    RU1285/18      GEMECHU SAMUEL MOGASA
10   RR0124/17      HAFSA SHAFI YASIN

Addis Ababa University | College of Natural Science | School of Information Science"""

# SLIDE 2: TABLE OF CONTENTS (Only for Computer Vision)
add_content_slide("Table of Contents", [
    "",
    "1.  What is Computer Vision?",
    "2.  History of Computer Vision",
    "3.  How Computer Vision Works (3-Step Process)",
    "4.  Types of Computer Vision Techniques (7 Types)",
    "5.  Real-World Applications of Computer Vision",
    "6.  Advantages (Pros) of Computer Vision",
    "7.  Disadvantages (Cons) of Computer Vision",
    "8.  Summary",
    "9.  References"
], font_size=20)

# SLIDE 3: What is Computer Vision?
add_content_slide("1. What is Computer Vision?", [
    "",
    "Definition:",
    "",
    "\"Computer vision is an interdisciplinary scientific field that",
    "deals with how computers can be made to gain a high-level",
    "understanding of digital images or videos.\"",
    "",
    "From an engineering perspective:",
    "• Seeks to automate tasks that the human visual system can do",
    "",
    "Key concept - 'Understanding' means:",
    "• Transforming visual images into descriptions of the world",
    "• Disentangling symbolic information from image data",
    "• Using models from geometry, physics, statistics, and learning theory"
], font_size=18)

# SLIDE 4: History of Computer Vision
add_content_slide("2. History of Computer Vision", [
    "",
    "1950s:",
    "• First neural networks for edge detection",
    "• Simple object sorting (circles and squares)",
    "",
    "1966:",
    "• MIT undergraduate summer project",
    "• Believed solvable in one summer!",
    "",
    "1970s:",
    "• First commercial use: Optical Character Recognition (OCR)",
    "• Used to interpret written text for the blind",
    "",
    "1990s:",
    "• Internet maturation → large image datasets available",
    "• Facial recognition programs flourished",
    "",
    "TODAY:",
    "• 50+ year old field still far from being solved!"
], font_size=17)

# SLIDE 5: How Computer Vision Works - Step 1
add_content_slide("3. How It Works - Step 1: Acquiring", [
    "",
    "STEP 1: ACQUIRING AN IMAGE",
    "",
    "• Images can be acquired in REAL-TIME",
    "",
    "Acquisition methods include:",
    "   - Video streams",
    "   - Digital photographs",
    "   - 3D technology",
    "",
    "• Even LARGE SETS of images can be acquired",
    "• Purpose: Get raw visual data for analysis"
], font_size=20)

# SLIDE 6: How It Works - Step 2 & 3
add_content_slide("3. How It Works - Steps 2 & 3", [
    "",
    "STEP 2: PROCESSING THE IMAGE",
    "",
    "• Deep learning models automate processing",
    "• Models are TRAINED with thousands of",
    "   labeled or pre-identified images",
    "• Training is essential for accuracy",
    "",
    "STEP 3: UNDERSTANDING THE IMAGE",
    "",
    "• The INTERPRETATIVE step",
    "• Where an object is IDENTIFIED or CLASSIFIED",
    "• Final output is meaningful information"
], font_size=20)

# SLIDE 7: 7 Types of Computer Vision - Part 1
add_content_slide("4. Types of CV Techniques - Part 1", [
    "",
    "1. IMAGE SEGMENTATION:",
    "   • Partitions an image into multiple regions or pieces",
    "   • Examines each region separately",
    "",
    "2. OBJECT DETECTION:",
    "   • Identifies a specific object in an image",
    "   • Advanced version recognizes MANY objects in one image",
    "   • Uses X,Y coordinates to create bounding boxes",
    "",
    "3. FACIAL RECOGNITION:",
    "   • Advanced type of object detection",
    "   • Not only recognizes a human face",
    "   • Identifies a SPECIFIC individual"
], font_size=17)

# SLIDE 8: 7 Types of Computer Vision - Part 2
add_content_slide("4. Types of CV Techniques - Part 2", [
    "",
    "4. EDGE DETECTION:",
    "   • Identifies the outside edge of an object or landscape",
    "   • Helps better identify what is in the image",
    "",
    "5. PATTERN DETECTION:",
    "   • Recognizes REPEATED shapes, colors, and visual indicators",
    "",
    "6. IMAGE CLASSIFICATION:",
    "   • Groups images into DIFFERENT CATEGORIES",
    "",
    "7. FEATURE MATCHING:",
    "   • Type of pattern detection",
    "   • Matches SIMILARITIES in images",
    "   • Helps CLASSIFY images"
], font_size=17)

# SLIDE 9: Important Note - Multiple Techniques
add_content_slide("Important Note", [
    "",
    "\"Simple applications of computer vision may only use",
    "ONE of these techniques...\"",
    "",
    "\"...but more advanced users, like computer vision for",
    "SELF-DRIVING CARS, rely on MULTIPLE techniques",
    "to accomplish their goal.\"",
    "",
    "",
    "Example - Self-driving car needs:",
    "   • Object detection (other cars, pedestrians)",
    "   • Edge detection (road lanes)",
    "   • Pattern detection (traffic signs)",
    "   • Image segmentation (different road elements)"
], font_size=20)

# SLIDE 10: Real-World Applications - Part 1
add_content_slide("5. Applications of CV - Part 1", [
    "",
    "Optical Character Recognition (OCR):",
    "   • Reading handwritten postal codes",
    "   • Automatic number plate recognition (ANPR)",
    "",
    "Machine Inspection:",
    "   • Rapid parts inspection for quality assurance",
    "   • Measuring tolerances on aircraft wings and auto body parts",
    "   • Detecting defects in steel castings using X-ray vision",
    "",
    "Retail:",
    "   • Object recognition for automated checkout lanes"
], font_size=17)

# SLIDE 11: Real-World Applications - Part 2
add_content_slide("5. Applications of CV - Part 2", [
    "",
    "Medical Imaging:",
    "   • Registering pre-operative and intra-operative imagery",
    "   • Studying brain morphology as people age",
    "",
    "Automotive Safety:",
    "   • Detecting unexpected obstacles like pedestrians",
    "   • Works where radar or lidar perform poorly",
    "",
    "Surveillance:",
    "   • Monitoring for intruders",
    "   • Analyzing highway traffic",
    "   • Monitoring pools for drowning victims",
    "",
    "Fingerprint Recognition & Biometrics:",
    "   • Automatic access authentication",
    "   • Forensic applications"
], font_size=17)

# SLIDE 12: Advantages (Pros) - Part 1
add_content_slide("6. Advantages of CV - Pros (Part 1)", [
    "",
    "✓ AUTOMATES HUMAN VISUAL TASKS",
    "   Seeks to automate what the human visual system can do",
    "",
    "✓ WORKS IN CONDITIONS HUMANS CAN'T",
    "   X-ray vision for detecting steel defects",
    "   Specialized illumination for precision measurement",
    "",
    "✓ HIGH PRECISION",
    "   Measures tolerances on aircraft wings",
    "   Quality assurance for auto body parts"
], font_size=19)

# SLIDE 13: Advantages (Pros) - Part 2
add_content_slide("6. Advantages of CV - Pros (Part 2)", [
    "",
    "✓ 24/7 CONTINUOUS OPERATION",
    "   Surveillance systems work around the clock",
    "   Never gets tired or loses attention",
    "",
    "✓ REAL-TIME PROCESSING",
    "   Can acquire and process images instantly",
    "   Critical for autonomous vehicles",
    "",
    "✓ SAVES LIVES",
    "   Pedestrian detection in cars prevents accidents",
    "   Drowning prevention in swimming pools",
    "",
    "✓ MULTIPLE TECHNIQUE INTEGRATION",
    "   Self-driving cars combine many techniques simultaneously"
], font_size=18)

# SLIDE 14: Disadvantages (Cons) - Part 1
add_content_slide("7. Disadvantages of CV - Cons (Part 1)", [
    "",
    "✗ STILL UNSOLVED AFTER 50+ YEARS",
    "   \"Still far from being solved\"",
    "   Originally believed solvable in one summer!",
    "",
    "✗ REQUIRES MASSIVE TRAINING DATA",
    "   Needs thousands of labeled or pre-identified images",
    "   Data collection is time-consuming and expensive",
    "",
    "✗ COMPUTATIONALLY EXPENSIVE",
    "   Requires deep learning models",
    "   Needs GPUs and specialized hardware"
], font_size=19)

# SLIDE 15: Disadvantages (Cons) - Part 2
add_content_slide("7. Disadvantages of CV - Cons (Part 2)", [
    "",
    "✗ LACKS HUMAN-LEVEL UNDERSTANDING",
    "   Cannot fully match human visual capabilities",
    "   Limited to what it was trained on",
    "",
    "✗ DEPENDENT ON TRAINING QUALITY",
    "   Simple applications may only use one technique",
    "   Poor training data = poor results",
    "",
    "✗ SENSOR AND LIGHTING DEPENDENT",
    "   Performance varies with image quality",
    "   Affected by lighting conditions, camera angles, and occlusion"
], font_size=18)

# SLIDE 16: Summary
add_content_slide("8. Summary", [
    "",
    "COMPUTER VISION:",
    "   Computers gaining high-level understanding of images and videos",
    "",
    "3-STEP PROCESS:",
    "   1. Acquire image (real-time, video, photo, 3D technology)",
    "   2. Process image (deep learning with trained models)",
    "   3. Understand image (identify and classify objects)",
    "",
    "7 MAIN TECHNIQUES:",
    "   Segmentation, Object Detection, Facial Recognition,",
    "   Edge Detection, Pattern Detection, Classification, Feature Matching",
    "",
    "PROS: Automation, 24/7 operation, high precision, life-saving",
    "",
    "CONS: Unsolved problem, needs massive data, computationally expensive"
], font_size=17)

# SLIDE 17: References (Only textbook, NO page numbers in content)
add_content_slide("9. References", [
    "",
    "Addis Ababa University, College of Natural Science,",
    "School of Information Science. (2019).",
    "Introduction to Emerging Technology.",
    "",
    "Chapter 7.6: Computer Vision",
    "",
    "(All information presented is based solely on",
    "the above textbook from Addis Ababa University)"
], font_size=20)

# SLIDE 18: Q&A
add_content_slide("10. Questions & Answers", [
    "",
    "",
    "",
    "              Thank you for your attention!",
    "",
    "",
    "",
    "              Any questions?",
    "",
    "",
    "              Based on:",
    "              Introduction to Emerging Technology (2019)",
    "              Addis Ababa University"
], font_size=24)

# Save
pr.save("Computer_Vision_Presentation_Final.pptx")
print("=" * 60)
print("✅ PRESENTATION CREATED SUCCESSFULLY")
print("=" * 60)
print(f"📊 Total slides: {len(pr.slides)}")
print("📚 Source: Textbook only (no page numbers shown in slide content)")
print("🎯 Topic: Computer Vision - How it works, Pros & Cons")
print("👥 Group members: 10 students listed on cover page")
print("=" * 60)
print("File saved as: Computer_Vision_Presentation_Final.pptx")